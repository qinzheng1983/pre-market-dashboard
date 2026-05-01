#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
长城汽车PPT生成脚本 - 2026年全球汽车行业"寒气"深度研究报告
基于python-pptx库，采用长城汽车品牌模板样式
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============ 长城汽车品牌色彩规范 ============
GWM_RED = RGBColor(230, 0, 18)  # #E60012 长城红
GWM_RED_DARK = RGBColor(196, 18, 48)  # #C41230 深红（警示）
GWM_GRAY = RGBColor(51, 51, 51)  # #333333 正文灰
GWM_GRAY_LIGHT = RGBColor(128, 128, 128)  # #808080 注释灰
GWM_BG_START = RGBColor(245, 245, 245)  # #F5F5F5 背景渐变起点
GWM_BG_END = RGBColor(255, 255, 255)  # #FFFFFF 背景渐变终点

# ============ PPT尺寸设置 ============
SLIDE_WIDTH = Inches(13.333)  # 16:9宽屏
SLIDE_HEIGHT = Inches(7.5)

class GWMPPTGenerator:
    """长城汽车风格PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        
    def add_title_bar(self, slide, title_text):
        """添加顶部红色条带标题"""
        # 红色条带背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_WIDTH, Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GWM_RED
        shape.line.fill.background()
        
        # 标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.15),
            Inches(10), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = '微软雅黑'
        
        # 长城汽车Logo占位（右上角）
        logo_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(0.15),
            Inches(1.5), Inches(0.5)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "长城汽车"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.RIGHT
        
    def add_page_number(self, slide, current, total):
        """添加页码"""
        page_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(7.1),
            Inches(1.5), Inches(0.3)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"第{current}页，共{total}页"
        p.font.size = Pt(10)
        p.font.color.rgb = GWM_GRAY_LIGHT
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.RIGHT
        
    def create_cover_slide(self):
        """创建封面页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 红色条带（比内页更宽）
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.5),
            SLIDE_WIDTH, Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GWM_RED
        shape.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.6),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "2026年全球汽车行业\"寒气\"深度研究报告"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = '微软雅黑'
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5),
            Inches(12), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "降本增效战略背景分析"
        p.font.size = Pt(20)
        p.font.color.rgb = GWM_GRAY
        p.font.name = '微软雅黑'
        
        # 日期和对象
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(12), Inches(0.5)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = "报告对象：跨国汽车集团CFO    |    2026年3月"
        p.font.size = Pt(14)
        p.font.color.rgb = GWM_GRAY_LIGHT
        p.font.name = '微软雅黑'
        
        # 页码
        self.add_page_number(slide, 1, 14)
        
    def create_summary_slide(self):
        """创建执行摘要页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "执行摘要 - 核心结论")
        
        # 关键指标卡片
        metrics = [
            ("中国批发销量", "-3%", "📉 首次负增长"),
            ("中国零售销量", "-7%", "📉 深度萎缩"),
            ("Q1销量预测", "-20%", "📉 断崖式下跌"),
            ("大众利润跌幅", "-53%", "📉 历史低位"),
            ("保时捷利润跌幅", "-92.7%", "📉 跌破预期"),
        ]
        
        y_pos = 1.2
        for i, (label, value, trend) in enumerate(metrics):
            x_pos = 0.3 + (i % 3) * 4.3
            if i == 3:
                y_pos = 3.5
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(4), Inches(2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 240, 240)
            card.line.color.rgb = GWM_RED
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_pos + 0.2),
                Inches(3.6), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = GWM_GRAY
            p.font.name = '微软雅黑'
            
            # 数值
            value_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_pos + 0.7),
                Inches(3.6), Inches(0.6)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = GWM_RED_DARK
            p.font.name = '微软雅黑'
            
            # 趋势
            trend_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_pos + 1.4),
                Inches(3.6), Inches(0.4)
            )
            tf = trend_box.text_frame
            p = tf.paragraphs[0]
            p.text = trend
            p.font.size = Pt(12)
            p.font.color.rgb = GWM_GRAY_LIGHT
            p.font.name = '微软雅黑'
        
        # 核心结论
        conclusion_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.8),
            Inches(12), Inches(0.8)
        )
        tf = conclusion_box.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 核心结论：行业从增量竞争正式进入存量博弈时代"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GWM_RED
        p.font.name = '微软雅黑'
        
        self.add_page_number(slide, 2, 14)
        
    def create_market_forecast_slide(self):
        """创建全球销量预测页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "全球销量预测 - 主要市场同步承压")
        
        # 表格数据
        headers = ["市场/指标", "预测值", "同比变化", "核心驱动因素"]
        data = [
            ["中国批发销量", "2,850万辆", "-3%", "购置税恢复、补贴退坡"],
            ["中国零售销量", "2,400万辆", "-7%", "消费者信心低迷"],
            ["中国Q1销量", "550万辆", "-20%", "政策切换、春节错期"],
            ["美国电动车", "110万辆", "-20%", "税收抵免到期"],
            ["欧洲电动车渗透率", "~18%", "持平", "碳标准放宽"],
            ["全球汽车增速", "~1.5%", "大幅放缓", "三大市场同步承压"],
        ]
        
        # 创建表格
        rows = len(data) + 1
        cols = len(headers)
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.3),
            Inches(12), Inches(5)
        ).table
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = GWM_RED
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.name = '微软雅黑'
            
        # 填充数据
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(11)
                paragraph.font.name = '微软雅黑'
                
                # 负增长用红色标注
                if col_idx == 2 and '-' in cell_text:
                    paragraph.font.color.rgb = GWM_RED_DARK
                    paragraph.font.bold = True
                    
        self.add_page_number(slide, 3, 14)
        
    def create_profit_collapse_slide(self):
        """创建龙头企业盈利崩塌页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "龙头企业盈利崩塌 - 量利双杀")
        
        # 标题说明
        desc_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1),
            Inches(12), Inches(0.4)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = "表2：2025-2026年汽车行业龙头企业盈利崩塌典型案例"
        p.font.size = Pt(12)
        p.font.color.rgb = GWM_GRAY_LIGHT
        p.font.name = '微软雅黑'
        
        # 表格
        headers = ["企业", "2025年关键财务指标", "同比变化", "核心压力来源"]
        data = [
            ["大众汽车集团", "营业利润85亿欧元，利润率3.5%", "-53%", "电动车转型、中国市场下滑、美国关税"],
            ["保时捷", "营业利润5.3亿欧元，回报率1.1%", "-92.7%", "中国市场暴跌26%、电动化战略调整"],
            ["博世集团", "销售额905亿欧元", "-1.2%", "电动化需求结构变化、欧洲能源成本"],
            ["本田汽车", "预计净亏损4,200-6,900亿日元", "由盈转亏", "电动化战略重组、资产减值2.5万亿日元"],
            ["中升控股", "预计亏损20亿元（去年盈利32亿）", "由盈转亏", "新车销售毛损扩大70%、金融佣金下降50%"],
        ]
        
        rows = len(data) + 1
        cols = len(headers)
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.6),
            Inches(12), Inches(5)
        ).table
        
        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = GWM_RED
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.name = '微软雅黑'
            
        # 数据
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.name = '微软雅黑'
                
                if col_idx == 2 and ('-' in cell_text or '亏损' in cell_text):
                    paragraph.font.color.rgb = GWM_RED_DARK
                    paragraph.font.bold = True
                    
        self.add_page_number(slide, 4, 14)
        
    def create_gdp_impact_slide(self):
        """创建宏观经济拖累测算页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "宏观经济拖累测算 - 汽车行业占GDP 3.2%")
        
        # 左侧：传导机制流程图
        flow_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(5.5), Inches(3)
        )
        tf = flow_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "GDP拖累传导机制："
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GWM_GRAY
        p.font.name = '微软雅黑'
        
        p = tf.add_paragraph()
        p.text = "\n汽车销量下滑"
        p.font.size = Pt(11)
        p.font.name = '微软雅黑'
        
        p = tf.add_paragraph()
        p.text = "↓ 直接效应：制造业增加值下降"
        p.font.size = Pt(11)
        p.font.color.rgb = GWM_RED
        p.font.name = '微软雅黑'
        
        p = tf.add_paragraph()
        p.text = "↓ 间接效应：产业链联动收缩"
        p.font.size = Pt(11)
        p.font.name = '微软雅黑'
        
        p = tf.add_paragraph()
        p.text = "↓ 诱发效应：收入减少→消费收缩"
        p.font.size = Pt(11)
        p.font.name = '微软雅黑'
        
        p = tf.add_paragraph()
        p.text = "\n= Q1拖累GDP增速 0.46个百分点"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GWM_RED_DARK
        p.font.name = '微软雅黑'
        
        # 右侧：对比表格
        comparison_box = slide.shapes.add_textbox(
            Inches(6.5), Inches(1.3),
            Inches(6), Inches(0.4)
        )
        tf = comparison_box.text_frame
        p = tf.paragraphs[0]
        p.text = "2018年 vs 2026年下行冲击对比"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GWM_GRAY
        p.font.name = '微软雅黑'
        
        # 对比表格
        headers = ["对比维度", "2018年", "2026年", "影响评估"]
        data = [
            ["驱动因素", "购置税退出（单一）", "购置税+补贴+竞争+外部（多重）", "2026年更复杂"],
            ["新能源市场", "占比<5%", "占比~50%", "2026年缺乏缓冲"],
            ["宏观环境", "中高速增长", "结构性增速下台阶", "2026年复苏基础薄弱"],
            ["政策空间", "较大", "有限", "2026年应对工具减少"],
        ]
        
        rows = len(data) + 1
        cols = len(headers)
        table = slide.shapes.add_table(
            rows, cols,
            Inches(6.5), Inches(1.8),
            Inches(6), Inches(4)
        ).table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = GWM_RED
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.name = '微软雅黑'
            
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.name = '微软雅黑'
                
        self.add_page_number(slide, 7, 14)
        
    def create_cost_reduction_slide(self):
        """创建降本增效四大抓手页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "降本增效四大关键抓手")
        
        # 四个抓手卡片
        levers = [
            {
                "title": "抓手1：供应链韧性重构",
                "items": ["• 多元化采购与供应商培育", "• 库存优化与风险管理", "• 物流网络重构"]
            },
            {
                "title": "抓手2：产能利用率提升",
                "items": ["• 需求驱动的生产计划", "• 柔性制造能力建设", "• 产能整合与退出"]
            },
            {
                "title": "抓手3：研发投入聚焦",
                "items": ["• 技术路线收敛", "• 研发效率提升", "• 商业化加速"]
            },
            {
                "title": "抓手4：组织效能提升",
                "items": ["• 组织架构扁平化", "• 数字化转型", "• 人员结构优化"]
            },
        ]
        
        positions = [(0.3, 1.2), (6.8, 1.2), (0.3, 4.2), (6.8, 4.2)]
        
        for idx, (lever, (x, y)) in enumerate(zip(levers, positions)):
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(6), Inches(2.7)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 248, 248)
            card.line.color.rgb = GWM_RED
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.2),
                Inches(5.6), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = lever["title"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = GWM_RED
            p.font.name = '微软雅黑'
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.8),
                Inches(5.6), Inches(1.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(lever["items"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(12)
                p.font.color.rgb = GWM_GRAY
                p.font.name = '微软雅黑'
                p.space_after = Pt(8)
                
        self.add_page_number(slide, 13, 14)
        
    def create_conclusion_slide(self):
        """创建结语页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.add_title_bar(slide, "结语与战略启示")
        
        # 行业周期判断
        period_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(12), Inches(1.5)
        )
        tf = period_box.text_frame
        
        p = tf.paragraphs[0]
        p.text = "行业周期判断："
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GWM_GRAY
        p.font.name = '微软雅黑'
        
        transitions = [
            "1. 从增量竞争 → 存量博弈",
            "2. 从政策驱动 → 市场驱动",
            "3. 从规模扩张 → 效率优先"
        ]
        
        for text in transitions:
            p = tf.add_paragraph()
            p.text = "\n" + text
            p.font.size = Pt(14)
            p.font.color.rgb = GWM_RED_DARK
            p.font.name = '微软雅黑'
            
        # 核心结论
        conclusion_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2),
            Inches(12), Inches(1.5)
        )
        tf = conclusion_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "\"寒气\"已不再是预警，而是现实。降本增效已从战略选项变为生存必需。"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GWM_RED
        p.font.name = '微软雅黑'
        
        # 历史经验
        history_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5),
            Inches(12), Inches(1.5)
        )
        tf = history_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "历史经验：行业寒冬也是格局重塑的窗口期——能够果断调整、快速执行、持续创新的企业，将在下一轮复苏中占据更有利的位置。"
        p.font.size = Pt(14)
        p.font.color.rgb = GWM_GRAY
        p.font.name = '微软雅黑'
        
        self.add_page_number(slide, 14, 14)
        
    def generate_full_ppt(self, output_path="gwm_auto_industry_report_2026.pptx"):
        """生成完整PPT"""
        print("🚗 开始生成长城汽车风格PPT...")
        
        # 创建各页
        self.create_cover_slide()
        print("✅ 封面页完成")
        
        self.create_summary_slide()
        print("✅ 执行摘要页完成")
        
        self.create_market_forecast_slide()
        print("✅ 全球销量预测页完成")
        
        self.create_profit_collapse_slide()
        print("✅ 盈利崩塌页完成")
        
        # 简化版：添加占位页提示
        for page_num in [5, 6, 8, 9, 10, 11, 12]:
            slide_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(slide_layout)
            self.add_title_bar(slide, f"第{page_num}页 - 详见完整报告")
            
            placeholder = slide.shapes.add_textbox(
                Inches(4), Inches(3.5),
                Inches(5), Inches(1)
            )
            tf = placeholder.text_frame
            p = tf.paragraphs[0]
            p.text = "（此页内容详见报告原文）"
            p.font.size = Pt(14)
            p.font.color.rgb = GWM_GRAY_LIGHT
            p.font.name = '微软雅黑'
            
            self.add_page_number(slide, page_num, 14)
            print(f"✅ 第{page_num}页占位完成")
        
        self.create_cost_reduction_slide()
        print("✅ 降本增效页完成")
        
        self.create_conclusion_slide()
        print("✅ 结语页完成")
        
        # 保存文件
        self.prs.save(output_path)
        print(f"\n🎉 PPT生成完成！保存至: {output_path}")
        print(f"📊 共生成 {len(self.prs.slides)} 页幻灯片")
        
        return output_path


# ============ 主程序 ============
if __name__ == "__main__":
    generator = GWMPPTGenerator()
    output_file = generator.generate_full_ppt()
    print(f"\n💡 使用提示：")
    print(f"   1. 安装依赖: pip install python-pptx")
    print(f"   2. 运行脚本: python3 gwm_ppt_generator.py")
    print(f"   3. 在PowerPoint中打开生成的文件")
    print(f"   4. 如需替换Logo，请修改代码中的Logo部分")
