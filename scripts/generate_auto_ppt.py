#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成长城汽车风格的汽车行业研究报告PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 长城汽车配色
GWM_RED = RGBColor(0xE6, 0x00, 0x12)
DARK_RED = RGBColor(0xC4, 0x12, 0x30)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部红色条带
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GWM_RED
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_RED
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年3月"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 页脚
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "跨国汽车集团CFO"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, is_table=False, table_data=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部红色条带
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GWM_RED
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if is_table and table_data:
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5)).table
        
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                cell.text_frame.paragraphs[0].font.size = Pt(11 if i > 0 else 12)
                cell.text_frame.paragraphs[0].font.bold = (i == 0)
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GWM_RED
                    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
    
    return slide

def main():
    prs = create_presentation()
    
    # 第1页：封面
    add_title_slide(prs, '2026年全球汽车行业"寒气"深度研究报告', "降本增效战略背景分析")
    
    # 第2页：执行摘要
    add_content_slide(prs, "执行摘要 - 核心结论", [
        "中国批发销量：-3%（首次负增长）",
        "中国零售销量：-7%（深度萎缩）",
        "Q1销量预测：-20%（断崖式下跌）",
        "大众利润跌幅：-53%（历史低位）",
        "保时捷利润跌幅：-92.7%（跌破预期）",
        "",
        "核心结论：行业从增量竞争正式进入存量博弈时代"
    ])
    
    # 第3页：全球销量预测
    add_content_slide(prs, "全球销量预测对比", [], is_table=True, table_data=[
        ["市场/指标", "预测值", "同比变化", "核心驱动因素"],
        ["中国批发销量", "2,850万辆", "-3%", "购置税恢复、补贴退坡"],
        ["中国零售销量", "2,400万辆", "-7%", "消费者信心低迷"],
        ["中国Q1销量", "550万辆", "-20%", "政策切换、春节错期"],
        ["美国电动车", "110万辆", "-20%", "税收抵免到期"],
        ["全球汽车增速", "1.5%", "大幅放缓", "三大市场同步承压"]
    ])
    
    # 第4页：龙头企业盈利崩塌
    add_content_slide(prs, "龙头企业盈利崩塌", [], is_table=True, table_data=[
        ["企业", "关键指标", "同比变化", "核心压力来源"],
        ["大众汽车", "营业利润85亿欧元，利润率3.5%", "-53%", "电动车转型、中国市场、美国关税"],
        ["保时捷", "营业利润5.3亿欧元，回报率1.1%", "-92.7%", "中国暴跌26%、电动化调整"],
        ["博世集团", "销售额905亿欧元", "-1.2%", "电动化需求结构变化"],
        ["本田汽车", "净亏损4,200-6,900亿日元", "由盈转亏", "电动化重组、减值2.5万亿日元"],
        ["中升控股", "预亏20亿元", "由盈转亏", "新车毛损扩大70%"]
    ])
    
    # 第5页：宏观经济预测
    add_content_slide(prs, "宏观经济环境 - 机构预测对比", [], is_table=True, table_data=[
        ["机构", "全球增速", "关键假设", "汽车行业隐含增速"],
        ["IMF", "3.3%", "美中有韧性，其他疲软", "1.5%-2.5%"],
        ["世界银行", "2.7% / 2.6%", "通胀粘性、贸易碎片化", "1.0%-2.0%"],
        ["联合国", "2.7%", "贸易增速2.2%", "0.5%-1.5%"],
        ["野村证券", "—", "三大市场同步承压", "1.5%"]
    ])
    
    # 第6页：政策退坡影响
    add_content_slide(prs, "中国市场政策退坡影响", [
        "2024-2025年                    2026年1月1日起",
        "   购置税免征  ---------------->  恢复至5%",
        "   以旧换新补贴  --------------->  规模缩减",
        "   插混续航门槛43km  ---------->  提升至100km",
        "   规模扩张导向  --------------->  质量提升导向",
        "",
        "影响测算：",
        "• 20万新能源车购置税增加：10,000元",
        "• 约40%插混车型失去优惠资格",
        "• 预计5-8家尾部车企退出市场"
    ])
    
    # 第7页：宏观经济拖累
    add_content_slide(prs, "宏观经济拖累测算", [
        "GDP拖累传导机制:",
        "汽车销量下滑 → 制造业增加值下降 → 产业链联动收缩",
        "       直接效应(0.46pp)      间接效应(乘数效应)      诱发效应(消费收缩)",
        "",
        "2018年 vs 2026年对比：",
        "• 2026年驱动因素更复杂（多重叠加）",
        "• 2026年缺乏新能源缓冲（占比~50%）",
        "• 2026年复苏基础薄弱，政策空间有限"
    ])
    
    # 第8页：国际市场分化
    add_content_slide(prs, "国际市场分化", [
        "美国市场 - \"电动车之冬\":",
        "  • 电动车销量预计下滑20%",
        "  • 税收抵免到期引发需求枯竭",
        "  • 混动车型份额回升",
        "  • 特朗普关税政策压力",
        "",
        "欧洲市场 - \"转型阵痛\":",
        "  • 碳排放标准放宽",
        "  • 本土品牌竞争力下滑",
        "  • 能源成本高企",
        "  • 电动车渗透放缓"
    ])
    
    # 第9页：产能过剩危机
    add_content_slide(prs, "产能过剩危机", [], is_table=True, table_data=[
        ["层级", "企业类型", "占比", "产能利用率"],
        ["高位运行", "比亚迪、特斯拉、理想", "30%", ">80%"],
        ["中等水平", "二线新势力、转型传统车企", "40%", "50%-80%"],
        ["低位闲置", "尾部新势力、落后产能", "30%", "<50%"]
    ])
    
    # 第10页：供应链危机
    add_content_slide(prs, "供应链危机 - 芯片短缺", [], is_table=True, table_data=[
        ["芯片类型", "传统燃油车", "智能电动车", "2026年预测", "价格涨幅"],
        ["DRAM", "0.5-1GB", "8-16GB", "16-32GB", "+150%-300%"],
        ["NAND Flash", "2-4GB", "64-128GB", "128-256GB", "+200%-400%"],
        ["HBM", "0", "8-16GB", "16-32GB", "+300%-500%"]
    ])
    
    # 第11页：裁员风暴
    add_content_slide(prs, "龙头企业裁员风暴", [], is_table=True, table_data=[
        ["企业", "裁员规模", "时间节点", "区域分布"],
        ["大众集团", "5万人", "2026年前", "多品牌、多区域"],
        ["博世", "1,200人", "2026年底前", "79%在德国"],
        ["保时捷", "3,900人", "2029年前", "渐进式"],
        ["法雷奥", "1,150人", "全球", "970人在欧洲"]
    ])
    
    # 第12页：风险矩阵
    add_content_slide(prs, "核心风险矩阵", [
        "需求端风险（最高优先级）：",
        "  ★ 国内需求萎缩：销量下滑3%-7%，Q1或暴跌20%",
        "  ★ 出口增长受阻：美欧关税升级",
        "  ★ 价格竞争激化：利润率压缩至3%-5%",
        "",
        "成本端风险（最高优先级）：",
        "  ★ 存储芯片短缺：满足率<50%，价格+300%-500%",
        "  ★ 电池材料波动：锂、钴、镍价格波动",
        "  ★ 汇率波动：美元强势",
        "",
        "政策端风险（最高优先级）：",
        "  ★ 新能源补贴退坡：购置税恢复5%",
        "  ★ 贸易壁垒升级：美欧对华关税",
        "  ★ 投资审查强化：关键技术限制"
    ])
    
    # 第13页：降本增效抓手
    add_content_slide(prs, "降本增效四大抓手", [
        "抓手1: 供应链韧性重构",
        "    • 多元化采购与供应商培育  • 库存优化与风险管理  • 物流网络重构",
        "",
        "抓手2: 产能利用率提升",
        "    • 需求驱动的生产计划  • 柔性制造能力建设  • 产能整合与退出",
        "",
        "抓手3: 研发投入聚焦",
        "    • 技术路线收敛  • 研发效率提升  • 商业化加速",
        "",
        "抓手4: 组织效能提升",
        "    • 组织架构扁平化  • 数字化转型  • 人员结构优化"
    ])
    
    # 第14页：结语
    add_content_slide(prs, "结语与战略启示", [
        "行业周期判断：",
        "    1. 从增量竞争 → 存量博弈",
        "    2. 从政策驱动 → 市场驱动", 
        "    3. 从规模扩张 → 效率优先",
        "",
        '核心结论："寒气"已不再是预警，而是现实。',
        "         降本增效已从战略选项变为生存必需。",
        "",
        "历史经验：",
        "    行业寒冬也是格局重塑的窗口期——",
        "    能够果断调整、快速执行、持续创新的企业，",
        "    将在下一轮复苏中占据更有利的位置。"
    ])
    
    # 保存文件
    output_dir = "/root/.openclaw/workspace/reports"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "2026汽车行业寒气研究报告.pptx")
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
